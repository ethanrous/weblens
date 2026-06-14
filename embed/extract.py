"""File-type dispatch for text extraction; returns 1-indexed (page, text) pairs."""
from __future__ import annotations

import os


class ExtractionError(Exception):
    """Raised when a file's content cannot be extracted into text."""


# A page's extracted text paired with its 1-indexed source page number.
Page = tuple[int, str]


_PLAINTEXT_MIMES = {
    "text/plain", "text/markdown", "text/csv", "text/yaml",
    "application/json", "application/x-yaml",
}

_PLAINTEXT_EXTS = {
    ".txt", ".md", ".csv", ".log", ".json", ".yaml", ".yml",
    ".go", ".py", ".js", ".ts", ".tsx", ".vue", ".rs", ".java",
    ".c", ".cpp", ".h", ".hpp", ".sh", ".rb", ".kt", ".swift",
}

_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".heic", ".tif", ".tiff", ".bmp"}


def extract_text(path: str, mime_hint: str | None = None) -> list[Page]:
    """Extract text from a file as a list of (page, text) pairs; raises ExtractionError for unsupported types."""
    ext = os.path.splitext(path)[1].lower()
    mime = (mime_hint or "").lower()

    if mime in _PLAINTEXT_MIMES or ext in _PLAINTEXT_EXTS:
        return [(1, _read_plaintext(path))]
    if ext == ".pdf" or mime == "application/pdf":
        return _extract_pdf(path)
    if ext == ".docx":
        return [(1, _extract_docx(path))]
    if ext == ".xlsx":
        return _extract_xlsx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if mime.startswith("image/") or ext in _IMAGE_EXTS:
        return [(1, _ocr_image(path))]

    raise ExtractionError(f"unsupported file type: ext={ext} mime={mime}")


def _read_plaintext(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as fh:
        return fh.read()


def _extract_pdf(path: str) -> list[Page]:
    import pypdf

    try:
        reader = pypdf.PdfReader(path)
    except Exception as e:
        raise ExtractionError(f"pdf parse failed: {e}") from e

    if reader.is_encrypted:
        raise ExtractionError("pdf is encrypted")

    pages: list[Page] = []
    for page_index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if not text:
            try:
                pil_img = page.to_image() if hasattr(page, "to_image") else None
            except Exception:
                pil_img = None
            if pil_img is not None:
                text = _ocr(pil_img)
        if text:
            pages.append((page_index, text))
    return pages


def _extract_docx(path: str) -> str:
    from docx import Document as DocxDocument
    doc = DocxDocument(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_xlsx(path: str) -> list[Page]:
    from openpyxl import load_workbook
    wb = load_workbook(filename=path, read_only=True, data_only=True)
    pages: list[Page] = []
    for sheet_index, sheet in enumerate(wb.worksheets, start=1):
        rows = (
            " ".join(str(c) for c in row if c is not None)
            for row in sheet.iter_rows(values_only=True)
        )
        text = "\n".join(r for r in rows if r)
        if text:
            pages.append((sheet_index, text))
    return pages


def _extract_pptx(path: str) -> list[Page]:
    from pptx import Presentation
    prs = Presentation(path)
    pages: list[Page] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        parts = [shape.text for shape in slide.shapes
                 if hasattr(shape, "text") and shape.text]
        text = "\n".join(parts)
        if text:
            pages.append((slide_index, text))
    return pages


# Tesseract per-word confidence (0-100) below which a word is discarded.
_MIN_WORD_CONF = 60


def _looks_legible(text: str) -> bool:
    """Reject OCR output that is mostly punctuation or stray single letters."""
    tokens = text.split()
    if not tokens:
        return False
    nonspace = [c for c in text if not c.isspace()]
    alpha_ratio = sum(c.isalnum() for c in nonspace) / len(nonspace)
    avg_token_len = sum(len(t) for t in tokens) / len(tokens)
    return alpha_ratio >= 0.7 and avg_token_len >= 2.5


def _ocr(img) -> str:
    import pytesseract
    data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
    words = [w for w, c in zip(data["text"], data["conf"])
             if w.strip() and float(c) >= _MIN_WORD_CONF]
    text = " ".join(words)
    return text if _looks_legible(text) else ""


def _ocr_image(path: str) -> str:
    from PIL import Image
    try:
        import pytesseract  # ensure a clear error before opening the image
    except ImportError as e:
        raise ExtractionError(f"OCR unavailable: {e}") from e
    return _ocr(Image.open(path))
