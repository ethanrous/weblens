"""Regenerate fixture files for test_extract.py. Run from the embed/ directory:

    python test_fixtures/_generate.py

Requires: pypdf, python-docx, openpyxl, python-pptx, Pillow.
"""
from __future__ import annotations

import io
import os
from PIL import Image, ImageDraw, ImageFont
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader, PdfWriter
from pypdf.generic import RectangleObject, NameObject, NumberObject

HERE = os.path.dirname(os.path.abspath(__file__))


def _minimal_pdf(text: str, out_path: str) -> None:
    """Write a 1-page PDF with the given literal text via raw PDF stream."""
    text_escaped = text.replace("(", "\\(").replace(")", "\\)")
    content = f"BT /F1 18 Tf 72 720 Td ({text_escaped}) Tj ET".encode("latin-1")
    length = len(content)
    pdf = (
        b"%PDF-1.4\n"
        b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj\n"
        b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj\n"
        b"3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >> endobj\n"
        + f"4 0 obj << /Length {length} >> stream\n".encode("latin-1")
        + content
        + b"\nendstream endobj\n"
        b"5 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj\n"
        b"xref\n0 6\n0000000000 65535 f \n"
    )
    # Build xref offsets dynamically.
    obj_offsets = []
    for marker in (b"1 0 obj", b"2 0 obj", b"3 0 obj", b"4 0 obj", b"5 0 obj"):
        obj_offsets.append(pdf.index(marker))
    xref_lines = b""
    for off in obj_offsets:
        xref_lines += f"{off:010d} 00000 n \n".encode("ascii")
    # Strip placeholder "xref\n0 6\n..." we appended and rebuild trailer.
    head = pdf[: pdf.index(b"xref\n")]
    xref_start = len(head)
    body = (
        head
        + b"xref\n"
        + f"0 6\n0000000000 65535 f \n".encode("ascii")
        + xref_lines
        + b"trailer << /Size 6 /Root 1 0 R >>\n"
        + f"startxref\n{xref_start}\n%%EOF\n".encode("ascii")
    )
    with open(out_path, "wb") as fh:
        fh.write(body)


def make_pdf():
    _minimal_pdf("semantic search test marker", os.path.join(HERE, "sample.pdf"))


def make_encrypted_pdf():
    src = os.path.join(HERE, "sample.pdf")
    if not os.path.exists(src):
        make_pdf()
    reader = PdfReader(src)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt("password123")
    with open(os.path.join(HERE, "encrypted.pdf"), "wb") as fh:
        writer.write(fh)


def make_docx():
    doc = Document()
    doc.add_paragraph("document marker phrase")
    doc.save(os.path.join(HERE, "sample.docx"))


def make_xlsx():
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "spreadsheet marker"
    wb.save(os.path.join(HERE, "sample.xlsx"))


def make_pptx():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tx_box.text_frame.text = "slide marker"
    prs.save(os.path.join(HERE, "sample.pptx"))


def make_png():
    img = Image.new("RGB", (400, 80), color="white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 32)
    except OSError:
        font = ImageFont.load_default()
    draw.text((10, 20), "image OCR marker", fill="black", font=font)
    img.save(os.path.join(HERE, "sample.png"))


def main():
    make_pdf()
    make_encrypted_pdf()
    make_docx()
    make_xlsx()
    make_pptx()
    make_png()
    print("fixtures generated in", HERE)


if __name__ == "__main__":
    main()
